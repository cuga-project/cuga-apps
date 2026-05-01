"""
Text extraction from heterogeneous file formats.

Supported:  PDF, PPTX, Markdown, plain text, MP3/MP4/WAV/M4A/WEBM/MOV (via Whisper).
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf":      "pdf",
    ".pptx":     "pptx",
    ".ppt":      "pptx",
    ".md":       "markdown",
    ".markdown": "markdown",
    ".txt":      "text",
    ".rst":      "text",
    ".mp3":      "audio",
    ".wav":      "audio",
    ".m4a":      "audio",
    ".mp4":      "video",
    ".webm":     "video",
    ".mov":      "video",
    ".avi":      "video",
}


def classify_file(filepath: str) -> Optional[str]:
    """Return file category string or None if unsupported."""
    return SUPPORTED_EXTENSIONS.get(Path(filepath).suffix.lower())


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def extract_pdf(filepath: str) -> str:
    import pdfplumber
    parts: list[str] = []
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text.strip())
    return "\n\n".join(parts)


def extract_pptx(filepath: str) -> str:
    from pptx import Presentation  # type: ignore

    prs = Presentation(filepath)
    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = [f"## Slide {idx}"]
        title_shape = slide.shapes.title
        if title_shape and title_shape.text.strip():
            lines.append(f"Title: {title_shape.text.strip()}")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text or shape is title_shape:
                continue
            lines.append(text)
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf and notes_tf.text.strip():
                lines.append(f"Notes: {notes_tf.text.strip()}")
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def extract_markdown(filepath: str) -> str:
    return Path(filepath).read_text(encoding="utf-8", errors="replace")


def extract_text(filepath: str) -> str:
    return Path(filepath).read_text(encoding="utf-8", errors="replace")


def transcribe_audio_video(filepath: str) -> str:
    """Transcribe audio/video using faster-whisper (base model, CPU)."""
    try:
        from faster_whisper import WhisperModel  # type: ignore
    except ImportError:
        raise ImportError(
            "faster-whisper is required for audio/video transcription.\n"
            "Install: pip install faster-whisper"
        )
    model = WhisperModel("base", device="cpu", compute_type="int8")
    segments, _info = model.transcribe(filepath, beam_size=5)
    return " ".join(seg.text for seg in segments).strip()


# ---------------------------------------------------------------------------
# Unified entry point
# ---------------------------------------------------------------------------

def extract(filepath: str) -> tuple[str, str]:
    """
    Extract text from a file.

    Returns:
        (file_type, extracted_text)

    Raises:
        ValueError  if the extension is not supported
        Exception   propagated from the underlying extractor
    """
    ftype = classify_file(filepath)
    if ftype is None:
        raise ValueError(f"Unsupported file type: {Path(filepath).suffix!r} ({filepath})")

    if ftype == "pdf":
        return ftype, extract_pdf(filepath)
    elif ftype == "pptx":
        return ftype, extract_pptx(filepath)
    elif ftype == "markdown":
        return ftype, extract_markdown(filepath)
    elif ftype == "text":
        return ftype, extract_text(filepath)
    elif ftype in ("audio", "video"):
        return ftype, transcribe_audio_video(filepath)

    raise ValueError(f"Unhandled type {ftype!r}")  # should not reach

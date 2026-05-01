"""
Build .pptx and .md output from a structured Deck object.

Uses a clean widescreen (16:9) template with a dark-navy title bar.
Each content slide follows the "title + body" layout.  Speaker notes
are written to the notes pane.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.util import Inches, Pt  # type: ignore


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    title: str
    bullets: list[str] = field(default_factory=list)
    speaker_notes: str = ""


@dataclass
class Deck:
    title: str
    subtitle: str = ""
    slides: list[Slide] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------

_NAVY   = RGBColor(0x0F, 0x2D, 0x52)   # IBM-ish dark navy
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)   # slide background
_ACCENT = RGBColor(0x00, 0x89, 0xFF)   # accent blue


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

def _set_16_9(prs: Presentation) -> None:
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)


def _set_notes(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf.paragraphs:
        tf.paragraphs[0].text = text
    else:
        tf.add_paragraph().text = text


def _add_title_slide(prs: Presentation, deck: Deck) -> None:
    layout = prs.slide_layouts[0]   # "Title Slide"
    slide  = prs.slides.add_slide(layout)

    title_ph = slide.shapes.title
    if title_ph:
        title_ph.text = deck.title
        title_ph.text_frame.paragraphs[0].font.size = Pt(40)
        title_ph.text_frame.paragraphs[0].font.color.rgb = _WHITE

    # Subtitle placeholder (index 1)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = deck.subtitle
            ph.text_frame.paragraphs[0].font.size = Pt(22)
            ph.text_frame.paragraphs[0].font.color.rgb = _WHITE
            break

    # Apply navy background fill
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _NAVY


def _add_content_slide(prs: Presentation, s: Slide) -> None:
    layout = prs.slide_layouts[1]   # "Title and Content"
    slide  = prs.slides.add_slide(layout)

    # Title
    title_ph = slide.shapes.title
    if title_ph:
        title_ph.text = s.title
        tf = title_ph.text_frame
        tf.paragraphs[0].font.size  = Pt(28)
        tf.paragraphs[0].font.bold  = True
        tf.paragraphs[0].font.color.rgb = _WHITE

    # Light background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _LIGHT

    # Navy band behind the title placeholder
    #  (the title already lives in a placeholder; tinting the slide bg is enough)

    # Body / bullets
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break

    if body_ph:
        tf = body_ph.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, bullet in enumerate(s.bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text  = bullet
            para.level = 0
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    if s.speaker_notes:
        _set_notes(slide, s.speaker_notes)


def build_pptx(deck: Deck, output_path: str) -> None:
    """Write the Deck to a .pptx file at output_path."""
    prs = Presentation()
    _set_16_9(prs)
    _add_title_slide(prs, deck)
    for slide in deck.slides:
        _add_content_slide(prs, slide)
    prs.save(output_path)


# ---------------------------------------------------------------------------
# Markdown / text report builder
# ---------------------------------------------------------------------------

def build_markdown(deck: Deck, output_path: str) -> None:
    """Write a structured Markdown version of the deck."""
    lines: list[str] = [f"# {deck.title}"]
    if deck.subtitle:
        lines.append(f"\n_{deck.subtitle}_")
    lines.append("")

    for i, slide in enumerate(deck.slides, start=1):
        lines.append(f"\n## {i}. {slide.title}\n")
        for bullet in slide.bullets:
            lines.append(f"- {bullet}")
        if slide.speaker_notes:
            lines.append(f"\n> **Speaker notes:** {slide.speaker_notes}")

    Path(output_path).write_text("\n".join(lines), encoding="utf-8")

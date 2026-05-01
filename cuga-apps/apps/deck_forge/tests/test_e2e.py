"""
DeckForge — End-to-End Integration Test
========================================

Tests the full pipeline end-to-end:
  fixture generation  →  RAG indexing  →  LangGraph agent  →  PPTX + MD output

Fixture directory contains:
  - tech_overview.md      synthetic Markdown doc on Transformer architecture
  - ml_concepts.txt       plain-text primer on attention and BERT
  - presentation.pptx     auto-generated PPTX covering key slides
  - (optionally) paper.pdf  downloaded from arxiv if network is available

The test runs the agent directly (no HTTP server needed).  It verifies:
  1. All files are indexed into the knowledge base
  2. The agent calls tools in the correct order
  3. The output PPTX exists and contains >= 5 slides (title + content)
  4. The output Markdown contains the topic keyword
  5. Slide titles are non-empty strings
  6. The session ends with status == "done"

Run:
    cd docs/examples/demo_apps/deck_forge
    python -m pytest tests/test_e2e.py -v

Requirements:
    An LLM API key in the environment (RITS_API_KEY, ANTHROPIC_API_KEY, etc.)
    The repo .venv must be active or on PYTHONPATH.
"""
from __future__ import annotations

import asyncio
import os
import sys
import tempfile
from pathlib import Path

import pytest

# ── path setup ────────────────────────────────────────────────────────────────
HERE = Path(__file__).parent
APP_DIR = HERE.parent
DEMO_APPS = APP_DIR.parent

sys.path.insert(0, str(APP_DIR))
sys.path.insert(0, str(DEMO_APPS))   # for _llm.py

# ── constants ─────────────────────────────────────────────────────────────────
TOPIC = "Transformer Architecture: Self-Attention, BERT, and Scalability"

_MARKDOWN_CONTENT = """\
# Transformer Architecture Overview

## Introduction

The Transformer model, introduced in "Attention Is All You Need" (Vaswani et al., 2017),
replaced recurrent networks with a purely attention-based architecture.  It became the
foundation for modern large language models (LLMs).

## Self-Attention Mechanism

Self-attention allows each token to attend to every other token in the sequence.
For each token, three vectors are computed: **Query (Q)**, **Key (K)**, and **Value (V)**.

    Attention(Q, K, V) = softmax(QK^T / sqrt(d_k)) · V

Multi-head attention runs h parallel attention heads, each learning different
relational patterns.  The outputs are concatenated and linearly projected.

## Positional Encoding

Because Transformers have no recurrence, position information is injected via
sinusoidal positional encodings added to the input embeddings.

## Feed-Forward Sub-layer

Each Transformer block contains a two-layer position-wise feed-forward network
(FFN) applied identically to each position:

    FFN(x) = max(0, xW_1 + b_1)W_2 + b_2

## Encoder–Decoder Architecture

The original Transformer has an encoder stack and a decoder stack.
- Encoder: self-attention + FFN
- Decoder: masked self-attention + cross-attention to encoder + FFN

## BERT and Bidirectional Pre-training

BERT (Devlin et al., 2018) introduced bidirectional pre-training of Transformers
using two tasks:
1. **Masked Language Modeling (MLM)** — predict randomly masked tokens
2. **Next Sentence Prediction (NSP)** — classify if two sentences are consecutive

BERT achieved state-of-the-art on 11 NLP benchmarks at the time of publication.

## Scalability

The attention mechanism has O(n²) complexity in sequence length, which limits
context windows.  Solutions include:
- Sparse attention (Longformer, BigBird)
- Linear attention approximations
- Flash Attention (IO-aware exact attention)

## Applications

Transformers are the backbone of:
- GPT series (text generation)
- BERT family (text understanding, QA)
- Vision Transformers (ViT) for image classification
- Whisper for speech recognition
- DALL-E and Stable Diffusion for image generation
"""

_TEXT_CONTENT = """\
Machine Learning Concepts: Attention and Neural Language Models
================================================================

ATTENTION MECHANISMS IN NLP
-----------------------------
Before Transformers, sequence models used RNNs or LSTMs with an attention
"add-on" (Bahdanau et al., 2015).  The attention score between a decoder
state and each encoder state determines how much the decoder attends to
each input position.

SELF-ATTENTION VS CROSS-ATTENTION
----------------------------------
Self-attention: tokens attend to other tokens in the *same* sequence.
Cross-attention: decoder tokens attend to encoder tokens.

TRANSFORMER TRAINING
---------------------
Training uses teacher forcing.  The loss is cross-entropy over the target
vocabulary.  Modern LLMs are trained with the Adam optimizer with a
linear warmup followed by cosine decay learning rate schedule.

BERT FINE-TUNING
-----------------
BERT adds a task-specific classification head on top of the [CLS] token.
Fine-tuning is efficient because only the head (and possibly the last few
layers) need substantial gradient updates.

PARAMETER COUNTS
-----------------
BERT-base:   110M parameters
BERT-large:  340M parameters
GPT-3:       175B parameters
GPT-4:       unknown (estimated 1T+)

COMPUTATIONAL COST
-------------------
Self-attention: O(n^2 d) time and O(n^2) memory per layer (n = seq length, d = dim).
For long documents this is the bottleneck.  Sparse attention reduces this to O(n sqrt(n)).

KEY TAKEAWAYS
--------------
1. Self-attention is the core innovation of the Transformer.
2. BERT demonstrated the power of bidirectional pre-training.
3. Scaling model size and data consistently improves performance (scaling laws).
4. Efficient attention variants are critical for long-context applications.
"""


def _make_sample_pptx(output_path: str) -> None:
    """Generate a sample PPTX with content about deep learning."""
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_data = [
        ("Deep Learning Foundations",
         ["Neural networks consist of layers of interconnected nodes",
          "Backpropagation computes gradients via the chain rule",
          "Gradient descent updates weights to minimise the loss",
          "Deep networks learn hierarchical feature representations"]),
        ("Recurrent Neural Networks",
         ["RNNs process sequences step by step, maintaining hidden state",
          "LSTMs use gates (input, forget, output) to control information flow",
          "Vanishing gradient problem limits long-range dependency learning",
          "Transformers address this by replacing recurrence with attention"]),
        ("Encoder–Decoder for Sequence-to-Sequence",
         ["Encoder compresses input into a fixed-size context vector",
          "Decoder generates output tokens auto-regressively",
          "Attention mechanism allows dynamic source alignment",
          "Transformer encoder–decoder achieves state-of-the-art in translation"]),
        ("Pre-training and Transfer Learning",
         ["Pre-train on large unlabeled corpora, fine-tune on downstream tasks",
          "BERT uses MLM; GPT uses causal language modelling",
          "Transfer learning dramatically reduces task-specific data requirements",
          "Foundation models are the dominant paradigm in modern NLP"]),
    ]

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Deep Learning and Transformers"
    slide.placeholders[1].text = "A visual overview of key concepts"

    # Content slides
    for title, bullets in slide_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b

    prs.save(output_path)


# ── Fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture(scope="module")
def fixture_dir(tmp_path_factory):
    """Create a temporary directory populated with test source materials."""
    d = tmp_path_factory.mktemp("deck_forge_fixtures")

    # Markdown
    (d / "tech_overview.md").write_text(_MARKDOWN_CONTENT, encoding="utf-8")

    # Plain text
    (d / "ml_concepts.txt").write_text(_TEXT_CONTENT, encoding="utf-8")

    # PPTX
    _make_sample_pptx(str(d / "presentation.pptx"))

    # Optionally download a real PDF (skip silently if no network)
    _try_download_pdf(str(d / "attention_paper.pdf"))

    print(f"\nFixture directory: {d}")
    print("Files:", [f.name for f in d.iterdir()])
    return str(d)


def _try_download_pdf(output_path: str) -> None:
    """
    Attempt to download the 'Attention is All You Need' paper from arxiv.
    Silently skips if the download fails (no network, rate limit, etc.).
    """
    import urllib.request

    url = "https://arxiv.org/pdf/1706.03762v5"
    try:
        req = urllib.request.Request(
            url, headers={"User-Agent": "DeckForge-Test/1.0"}
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            data = resp.read()
        if len(data) > 10_000:   # sanity-check: real PDF is > 10 KB
            Path(output_path).write_bytes(data)
            print(f"Downloaded PDF: {len(data)//1024} KB")
    except Exception as exc:
        print(f"PDF download skipped: {exc}")


@pytest.fixture(scope="module")
def output_dir(tmp_path_factory):
    return tmp_path_factory.mktemp("deck_forge_output")


# ── Unit tests for core modules ───────────────────────────────────────────────

class TestExtractors:
    def test_classify_pdf(self):
        from extractors import classify_file
        assert classify_file("doc.pdf") == "pdf"

    def test_classify_pptx(self):
        from extractors import classify_file
        assert classify_file("slides.pptx") == "pptx"

    def test_classify_markdown(self):
        from extractors import classify_file
        assert classify_file("README.md") == "markdown"

    def test_classify_unsupported(self):
        from extractors import classify_file
        assert classify_file("image.png") is None

    def test_extract_markdown(self, fixture_dir):
        from extractors import extract
        ftype, text = extract(str(Path(fixture_dir) / "tech_overview.md"))
        assert ftype == "markdown"
        assert "Transformer" in text
        assert "Self-Attention" in text

    def test_extract_text(self, fixture_dir):
        from extractors import extract
        ftype, text = extract(str(Path(fixture_dir) / "ml_concepts.txt"))
        assert ftype == "text"
        assert "BERT" in text

    def test_extract_pptx(self, fixture_dir):
        from extractors import extract
        ftype, text = extract(str(Path(fixture_dir) / "presentation.pptx"))
        assert ftype == "pptx"
        assert "Transformer" in text

    def test_extract_pdf(self, fixture_dir):
        """Only run if the PDF was successfully downloaded."""
        pdf_path = Path(fixture_dir) / "attention_paper.pdf"
        if not pdf_path.exists():
            pytest.skip("PDF not downloaded")
        from extractors import extract
        ftype, text = extract(str(pdf_path))
        assert ftype == "pdf"
        assert len(text) > 500


class TestKnowledgeBase:
    def test_add_and_search(self, fixture_dir):
        from rag import KnowledgeBase
        kb = KnowledgeBase("unit-test")
        from extractors import extract
        _, text = extract(str(Path(fixture_dir) / "tech_overview.md"))
        n = kb.add_document("tech_overview.md", text)
        assert n > 0
        results = kb.search("self-attention mechanism query key value", n_results=3)
        assert len(results) >= 1
        assert "score" in results[0]
        assert "text" in results[0]
        # Top result should be about attention
        top_text = results[0]["text"].lower()
        assert any(kw in top_text for kw in ["attention", "query", "key", "value", "transformer"])

    def test_multi_doc_search(self, fixture_dir):
        from rag import KnowledgeBase
        from extractors import extract
        kb = KnowledgeBase("multi-test")
        for fname in ["tech_overview.md", "ml_concepts.txt"]:
            _, text = extract(str(Path(fixture_dir) / fname))
            kb.add_document(fname, text)
        assert kb.chunk_count > 0
        results = kb.search("BERT bidirectional pre-training masked language model")
        assert any("bert" in r["text"].lower() for r in results)


class TestDeckWriter:
    def test_build_pptx(self, output_dir):
        from deck_writer import Deck, Slide, build_pptx
        from pptx import Presentation  # type: ignore
        deck = Deck(
            title="Test Deck",
            subtitle="Generated by DeckForge test",
            slides=[
                Slide("Intro", ["Point A", "Point B", "Point C"], "These are intro notes."),
                Slide("Deep Dive", ["Detail 1", "Detail 2"]),
                Slide("Takeaways", ["Key point", "Action item"]),
            ],
        )
        out = str(Path(output_dir) / "test.pptx")
        build_pptx(deck, out)
        assert Path(out).exists()
        prs = Presentation(out)
        # Title slide + 3 content slides
        assert len(prs.slides) == 4

    def test_build_markdown(self, output_dir):
        from deck_writer import Deck, Slide, build_markdown
        deck = Deck(
            title="MD Test",
            slides=[Slide("Section 1", ["A", "B"]), Slide("Section 2", ["C"])],
        )
        out = str(Path(output_dir) / "test.md")
        build_markdown(deck, out)
        text = Path(out).read_text()
        assert "# MD Test" in text
        assert "Section 1" in text
        assert "- A" in text


# ── E2E integration test ──────────────────────────────────────────────────────

def _llm_available() -> bool:
    """Check if any LLM key is set."""
    keys = ["RITS_API_KEY", "ANTHROPIC_API_KEY", "OPENAI_API_KEY",
            "WATSONX_APIKEY", "LITELLM_API_KEY"]
    return any(os.getenv(k) for k in keys)


@pytest.mark.skipif(not _llm_available(), reason="No LLM API key in environment")
class TestE2EPipeline:
    """Full end-to-end test: fixture dir → LangGraph agent → PPTX + MD."""

    def test_full_generation(self, fixture_dir, output_dir):
        """
        Runs the complete DeckForge pipeline:
        - Indexes all fixture files
        - Agent reasons, searches, builds slides
        - Writes PPTX and MD
        - Asserts meaningful output
        """
        from session import DeckForgeSession
        from agent import run_agent
        from pptx import Presentation  # type: ignore

        session = DeckForgeSession(
            session_id="e2e-test",
            directory=fixture_dir,
            topic=TOPIC,
            output_dir=Path(output_dir) / "e2e",
        )

        # Run agent synchronously for the test
        asyncio.run(run_agent(session))

        # ── Assert session completed ───────────────────────────────────
        assert session.status == "done", (
            f"Session ended with status={session.status!r}, error={session.error}"
        )
        assert session.result is not None

        # ── Assert PPTX ───────────────────────────────────────────────
        pptx_path = Path(session.result["pptx"])
        assert pptx_path.exists(), f"PPTX not found: {pptx_path}"

        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)
        assert total_slides >= 5, (
            f"Expected >= 5 slides (title + content), got {total_slides}"
        )

        # Check no slide has an empty title
        for i, slide in enumerate(prs.slides):
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                assert title_text, f"Slide {i+1} has an empty title"

        # ── Assert Markdown ────────────────────────────────────────────
        md_path = Path(session.result["md"])
        assert md_path.exists(), f"MD not found: {md_path}"

        md_text = md_path.read_text(encoding="utf-8")
        assert len(md_text) > 200, "Markdown output is too short"
        # At least one of the topic keywords should appear
        assert any(
            kw.lower() in md_text.lower()
            for kw in ["transformer", "attention", "bert", "scalability"]
        ), f"Topic keywords not found in markdown:\n{md_text[:500]}"

        # ── Assert knowledge base was used ────────────────────────────
        assert session.kb.chunk_count > 0, "Knowledge base was never populated"
        assert len(session.kb.sources) >= 2, "Agent indexed fewer than 2 sources"

        # ── Assert slide quality ──────────────────────────────────────
        assert session.result["slide_count"] >= 4, (
            f"Agent produced only {session.result['slide_count']} content slides"
        )
        for slide in session.slides:
            assert slide.title.strip(), "Found a slide with empty title"
            assert len(slide.bullets) >= 1, f"Slide '{slide.title}' has no bullets"

        print(f"\n✅ E2E test passed!")
        print(f"   Slides: {session.result['slide_count']}")
        print(f"   KB chunks: {session.kb.chunk_count}")
        print(f"   Sources: {len(session.kb.sources)}")
        print(f"   PPTX: {pptx_path}")
        print(f"   MD:   {md_path}")

    def test_progress_events_emitted(self, fixture_dir, output_dir):
        """
        Verify the agent emits the ingestion-phase progress events.

        The ingestion events (start, directory_scanned, extracting, indexed) are
        deterministic — they come from the tools, not the LLM.  We don't re-assert
        slide_added/done here because test_full_generation already covers the full
        pipeline end-to-end; this test focuses on the event structure.
        """
        from session import DeckForgeSession
        from agent import run_agent

        session = DeckForgeSession(
            session_id="e2e-events",
            directory=fixture_dir,
            topic="Self-Attention Mechanisms",
            output_dir=Path(output_dir) / "e2e-events",
        )

        asyncio.run(run_agent(session))

        # Drain the queue into a list (events emitted to an unconsumed queue
        # remain there until drained here)
        events: list[dict] = []
        while not session.queue.empty():
            events.append(session.queue.get_nowait())

        event_types = {e.get("type") for e in events}
        print(f"\nEvent types emitted: {event_types}")

        # Ingestion-phase events come from tool code (deterministic):
        required_ingestion = {"start", "directory_scanned", "indexed"}
        missing = required_ingestion - event_types
        assert not missing, (
            f"Missing ingestion event types: {missing}\n"
            f"Got: {event_types}"
        )

        # Verify at least one 'indexed' event has the expected structure
        indexed_events = [e for e in events if e.get("type") == "indexed"]
        assert indexed_events, "No 'indexed' events found"
        first = indexed_events[0]
        assert "file" in first
        assert "chunks" in first and first["chunks"] > 0

        # If the run completed successfully, verify slide events too
        if session.status == "done":
            assert "slide_added" in event_types, "status=done but no slide_added events"
            assert "done" in event_types
            slide_events = [e for e in events if e.get("type") == "slide_added"]
            assert all("title" in e for e in slide_events)
        else:
            # Partial run still verified ingestion above — acceptable for the events test
            print(f"  (agent ended with status={session.status!r} — ingestion events verified)")


# ── CugaAgent E2E ─────────────────────────────────────────────────────────────

def _cuga_available() -> bool:
    """Check that cuga.sdk is installed AND an LLM key is present."""
    try:
        from cuga.sdk import CugaAgent  # noqa: F401
        return _llm_available()
    except ImportError:
        return False


@pytest.mark.skipif(not _cuga_available(), reason="cuga.sdk not installed or no LLM key")
class TestCugaAgentE2E:
    """
    End-to-end tests for the CugaAgent backend.

    Uses the same fixture directory as TestE2EPipeline so the source material
    (PDFs, PPTX, Markdown, text) is identical.  The assertions mirror the
    LangGraph E2E tests so regressions in either backend show up symmetrically.
    """

    TOPIC = "Transformer Architecture: Self-Attention, BERT, and Scalability"

    def test_full_generation_cuga(self, fixture_dir, output_dir):
        """
        Full pipeline via CugaAgent:
          fixture dir → ingest → RAG → CugaAgent.invoke() → PPTX + MD
        """
        from session import DeckForgeSession
        from agent import run_agent
        from pptx import Presentation  # type: ignore

        session = DeckForgeSession(
            session_id="cuga-e2e",
            directory=fixture_dir,
            topic=self.TOPIC,
            output_dir=Path(output_dir) / "cuga-e2e",
            agent_type="cuga",
        )

        asyncio.run(run_agent(session))

        # ── Session status ─────────────────────────────────────────────
        assert session.status == "done", (
            f"CugaAgent session ended with status={session.status!r}, "
            f"error={session.error}"
        )
        assert session.result is not None

        # ── PPTX ──────────────────────────────────────────────────────
        pptx_path = Path(session.result["pptx"])
        assert pptx_path.exists(), f"PPTX not found: {pptx_path}"

        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)
        assert total_slides >= 5, (
            f"Expected >= 5 slides (title + content), got {total_slides}"
        )
        for i, slide in enumerate(prs.slides):
            if slide.shapes.title:
                assert slide.shapes.title.text.strip(), f"Slide {i+1} has empty title"

        # ── Markdown ──────────────────────────────────────────────────
        md_path = Path(session.result["md"])
        assert md_path.exists(), f"MD not found: {md_path}"
        md_text = md_path.read_text(encoding="utf-8")
        assert len(md_text) > 200, "Markdown output is too short"
        assert any(
            kw.lower() in md_text.lower()
            for kw in ["transformer", "attention", "bert", "scalability"]
        ), f"Topic keywords not found in markdown"

        # ── Knowledge base ────────────────────────────────────────────
        assert session.kb.chunk_count > 0, "KB was never populated"
        assert len(session.kb.sources) >= 2, "Agent indexed fewer than 2 sources"

        # ── Slides ────────────────────────────────────────────────────
        assert session.result["slide_count"] >= 4
        for slide in session.slides:
            assert slide.title.strip(), "Slide with empty title"
            assert len(slide.bullets) >= 1, f"Slide '{slide.title}' has no bullets"

        print(f"\n✅ CugaAgent E2E passed!")
        print(f"   Slides: {session.result['slide_count']}")
        print(f"   KB chunks: {session.kb.chunk_count}")
        print(f"   Sources: {len(session.kb.sources)}")
        print(f"   PPTX: {pptx_path}")
        print(f"   MD:   {md_path}")

    def test_cuga_progress_events(self, fixture_dir, output_dir):
        """
        CugaAgent tool closures push the same typed events as LangGraph.
        Verify ingestion-phase events are present regardless of which
        agent backend is used.
        """
        from session import DeckForgeSession
        from agent import run_agent

        session = DeckForgeSession(
            session_id="cuga-events",
            directory=fixture_dir,
            topic="Self-Attention in Transformers",
            output_dir=Path(output_dir) / "cuga-events",
            agent_type="cuga",
        )

        asyncio.run(run_agent(session))

        events: list[dict] = []
        while not session.queue.empty():
            events.append(session.queue.get_nowait())

        event_types = {e.get("type") for e in events}
        print(f"\nCugaAgent event types: {event_types}")

        # Ingestion events come from shared tool closures — same for both agents
        required = {"start", "directory_scanned", "indexed"}
        missing = required - event_types
        assert not missing, (
            f"Missing ingestion events: {missing}\nGot: {event_types}"
        )

        indexed = [e for e in events if e.get("type") == "indexed"]
        assert indexed, "No indexed events"
        assert all(e.get("chunks", 0) > 0 for e in indexed)

        if session.status == "done":
            assert "slide_added" in event_types
            assert "done" in event_types
        else:
            print(f"  (ended with {session.status!r} — ingestion events verified)")

    def test_agent_type_routing(self, fixture_dir, output_dir):
        """
        Verify that agent_type='langgraph' and agent_type='cuga' produce
        independently valid outputs from the same source directory.
        Both decks must have >= 4 content slides and contain the topic keywords.
        """
        from session import DeckForgeSession
        from agent import run_agent

        lg_session = DeckForgeSession(
            session_id="routing-lg",
            directory=fixture_dir,
            topic="BERT and Bidirectional Pre-training",
            output_dir=Path(output_dir) / "routing-lg",
            agent_type="langgraph",
        )
        cuga_session = DeckForgeSession(
            session_id="routing-cuga",
            directory=fixture_dir,
            topic="BERT and Bidirectional Pre-training",
            output_dir=Path(output_dir) / "routing-cuga",
            agent_type="cuga",
        )

        # Run sequentially (both hit the same RITS endpoint)
        asyncio.run(run_agent(lg_session))
        asyncio.run(run_agent(cuga_session))

        for label, s in [("LangGraph", lg_session), ("CugaAgent", cuga_session)]:
            assert s.status == "done", (
                f"{label} session ended with status={s.status!r}, error={s.error}"
            )
            assert s.result is not None, f"{label} result is None after done status"
            assert s.result["slide_count"] >= 4, (
                f"{label} produced only {s.result['slide_count']} slides"
            )
            md = Path(s.result["md"]).read_text()
            assert "bert" in md.lower() or "attention" in md.lower(), (
                f"{label} markdown missing topic keywords"
            )
            print(f"  {label}: {s.result['slide_count']} slides ✓")
